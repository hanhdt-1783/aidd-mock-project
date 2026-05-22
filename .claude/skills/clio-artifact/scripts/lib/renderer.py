from __future__ import annotations

import base64
import copy
import os
import tempfile
from pathlib import Path
from typing import Dict
from typing import List
from typing import Optional

from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt

from lib.schemas import ContentItem
from lib.schemas import ContentType
from lib.schemas import SlideContent
from lib.templates import SVNProposalTemplate



class PPTXRenderer:
    """Renderer for filling content into existing template slides"""

    # Maximum number of data rows (excluding header) allowed per slide table.
    # If a table exceeds this limit, the content is split across multiple slides.
    MAX_TABLE_ROWS_PER_SLIDE = 8

    def __init__(self):
        self.template_config = None
        self._output_dir: Optional[Path] = None

    def render(
        self,
        slides: List[SlideContent],
        template: str,
        output_name: str,
        output_dir: Optional[str] = None,
    ) -> Path:
        """
        Fill content into existing template slides

        Args:
            slides: List of slide content to fill
            template: Template PPTX file path (e.g., '[SVN] Proposal Menu.pptx')
            output_name: Output file name (without .pptx extension)
            output_dir: Deprecated/ignored. Output directory is resolved from env `SLIDE_GENERATOR__OUTPUTS_PATH`.

        Returns:
            Path to generated PPTX file
        """
        print(f'Filling content into template slides: {output_name}')

        # Resolve output_dir once so image path resolution can use it
        self._output_dir = self._resolve_output_dir(output_dir)

        # Auto-detect template config from template name
        self._detect_template_config(template)

        # Load template presentation
        prs = self._load_template(template)

        # Fill content into existing slides
        self._fill_existing_slides(prs, slides)

        # Save presentation
        output_file = self._save_presentation(prs, output_name, output_dir)

        return output_file

    @staticmethod
    def _resolve_output_dir(output_dir: Optional[str]) -> Path:
        """Resolve output directory path.

        Priority: explicit output_dir > env SLIDE_GENERATOR__OUTPUTS_PATH > CWD/outputs
        """
        if output_dir:
            return Path(output_dir)
        env_outputs = os.getenv('SLIDE_GENERATOR__OUTPUTS_PATH')
        if env_outputs:
            return Path(env_outputs).expanduser()
        return Path.cwd() / 'outputs'

    def _detect_template_config(self, template: str):
        """Auto-detect and load template config from template name"""
        if 'svn' in template.lower():
            self.template_config = SVNProposalTemplate()
            print(f'Auto-detected SVN template config: {len(self.template_config.slide_configs)} slides configured')
        else:
            self.template_config = None
            print('No specific template config detected, using default rendering')

    def _load_template(self, template: str) -> Presentation:
        """Load template presentation — supports absolute paths and relative names."""
        template_path = Path(template)
        if template_path.is_absolute() and template_path.exists():
            prs = Presentation(str(template_path))
            print(f'Loaded template: {template_path.name} ({len(prs.slides)} slides)')
            return prs
        rel_path = Path(__file__).parent / 'templates' / template
        if rel_path.exists():
            prs = Presentation(str(rel_path))
            print(f'Loaded template: {template} from {rel_path} ({len(prs.slides)} slides)')
            return prs
        print(f'Template not found: {template} at {template_path} or {rel_path}')
        return Presentation()

    def _save_presentation(self, prs: Presentation, output_name: str, output_dir: Optional[str] = None) -> Path:
        """Save presentation to output directory.

        Priority: explicit output_dir > env SLIDE_GENERATOR__OUTPUTS_PATH > CWD/outputs
        """
        output_path = self._output_dir or self._resolve_output_dir(output_dir)

        try:
            output_path.mkdir(parents=True, exist_ok=True)
        except PermissionError as e:
            # Make error clearer for API layer
            raise PermissionError(f'Permission denied creating outputs directory: {output_path}') from e
        output_file = output_path / f'{output_name}.pptx'

        prs.save(str(output_file))
        print(f'Saved PPTX to: {output_file}')

        return output_file

    def _fill_existing_slides(self, prs: Presentation, slides_content: List[SlideContent]):
        """Fill content into existing slides"""
        print(f'Filling content into {len(slides_content)} existing slides')

        # Sort by slide number so that the offset tracking works correctly when
        # overflow slides are inserted in the middle of the presentation.
        sorted_contents = sorted(slides_content, key=lambda s: s.slide_number)

        # Tracks how many extra slides have been inserted before the current one.
        slide_offset = 0

        for slide_content in sorted_contents:
            slide_num = slide_content.slide_number

            # Check if slide is protected
            if self.template_config and not self.template_config.should_fill_slide(slide_num):
                print(f'Slide {slide_num} is protected - skipping')
                continue

            # Adjusted index accounts for any extra slides inserted before this slide.
            adjusted_index = slide_num - 1 + slide_offset

            # Check if slide exists
            if adjusted_index < len(prs.slides):
                slide = prs.slides[adjusted_index]
                print(f'Filling content into slide {slide_num} (template index: {adjusted_index})')

                # Fill content; if a table overflows it will insert extra slides and
                # return the count so we can update the offset.
                extra = self._fill_with_overflow_handling(prs, slide, slide_content, adjusted_index)
                slide_offset += extra
            else:
                print(
                    f'Slide {slide_num} not found in template '
                    f'(total slides: {len(prs.slides)}, adjusted index: {adjusted_index})',
                )

    # ------------------------------------------------------------------
    # Table-overflow helpers
    # ------------------------------------------------------------------

    def _fill_with_overflow_handling(
        self,
        prs: Presentation,
        slide,
        slide_content: SlideContent,
        slide_index: int,
    ) -> int:
        """Fill a slide, splitting into multiple slides when a table exceeds
        MAX_TABLE_ROWS_PER_SLIDE data rows.

        Returns the number of *extra* slides inserted (0 when no overflow).
        """
        shape_contents = getattr(slide_content, 'shape_contents', {}) or {}

        # Find the first table-type shape content for this slide
        table_key, table_data = self._find_table_content(slide_content, shape_contents)

        if not table_key:
            # No table present – fill normally
            self._fill_slide_content(slide, slide_content)
            return 0

        # Split the table into row-capped chunks
        chunks = self._split_table_data(table_data, self.MAX_TABLE_ROWS_PER_SLIDE)

        if len(chunks) <= 1:
            # Fits on one slide – fill normally
            self._fill_slide_content(slide, slide_content)
            return 0

        print(
            f'  - Slide {slide_content.slide_number}: table overflow '
            f'({len(chunks)} page(s), max {self.MAX_TABLE_ROWS_PER_SLIDE} rows/page)',
        )

        # Fill the original slide with the first chunk
        first_content = self._build_chunked_slide_content(
            slide_content, table_key, chunks[0], is_first=True,
        )
        self._fill_slide_content(slide, first_content)

        # Create one continuation slide per remaining chunk
        extra_slides = 0
        current_index = slide_index
        for i, chunk in enumerate(chunks[1:], 1):
            new_slide = self._duplicate_slide(prs, current_index)
            extra_slides += 1
            current_index += 1

            cont_content = self._build_chunked_slide_content(
                slide_content, table_key, chunk, is_first=False,
            )
            self._fill_slide_content(new_slide, cont_content)
            print(
                f'  - Added continuation slide {i}/{len(chunks) - 1} '
                f'at template position {current_index + 1}',
            )

        return extra_slides

    def _find_table_content(
        self,
        slide_content: SlideContent,
        shape_contents: dict,
    ):
        """Return (content_key, markdown_table_string) for the first table-type
        shape found in this slide's config.  Returns (None, None) if none found.
        """
        config = (
            self.template_config.get_slide_config(slide_content.slide_number)
            if self.template_config
            else None
        )
        if not config or not config.shape_targets:
            return None, None

        for target in config.shape_targets:
            content = shape_contents.get(target.content_key, '')
            if content and self._is_table_data(content):
                return target.content_key, content

        return None, None

    def _split_table_data(self, table_data: str, max_rows: int) -> List[str]:
        """Split a markdown table string into chunks of at most *max_rows* data rows.

        Each chunk keeps the original header row(s) and the separator row so that
        each chunk is a valid, self-contained markdown table.

        Returns a list of markdown table strings (length 1 when no split needed).
        """
        lines = [line for line in table_data.strip().split('\n') if line.strip()]
        if not lines:
            return [table_data]

        # Separate header / separator lines from data lines
        header_lines: List[str] = []
        data_lines: List[str] = []
        separator_found = False

        for line in lines:
            stripped = line.strip()
            is_separator = (
                not separator_found
                and '|' in stripped
                and all(c in '|-: \t' for c in stripped)
            )
            if is_separator:
                header_lines.append(line)
                separator_found = True
            elif not separator_found:
                header_lines.append(line)
            else:
                data_lines.append(line)

        if len(data_lines) <= max_rows:
            return [table_data]

        header_block = '\n'.join(header_lines)
        chunks: List[str] = []
        for i in range(0, len(data_lines), max_rows):
            row_block = '\n'.join(data_lines[i : i + max_rows])
            chunks.append(f'{header_block}\n{row_block}')

        return chunks

    def _build_chunked_slide_content(
        self,
        original: SlideContent,
        table_key: str,
        table_chunk: str,
        is_first: bool,
    ) -> SlideContent:
        """Return a new SlideContent with *table_key* replaced by *table_chunk*.

        For continuation slides (``is_first=False``) all non-table text shapes
        are cleared so description text is not repeated across slides.
        """
        new_shape_contents = dict(getattr(original, 'shape_contents', {}) or {})
        new_shape_contents[table_key] = table_chunk

        if not is_first:
            config = (
                self.template_config.get_slide_config(original.slide_number)
                if self.template_config
                else None
            )
            if config and config.shape_targets:
                for target in config.shape_targets:
                    key = target.content_key
                    if key != table_key and key in new_shape_contents:
                        if not self._is_table_data(new_shape_contents[key]):
                            new_shape_contents[key] = ''  # Clear description text

        return SlideContent(
            slide_number=original.slide_number,
            layout=original.layout,
            title=original.title,
            subtitle=original.subtitle,
            content=original.content,
            shape_contents=new_shape_contents,
        )

    def _duplicate_slide(self, prs: Presentation, slide_index: int):
        """Duplicate the slide at *slide_index* and insert the copy immediately
        after it in the presentation.

        Returns the new slide object (at position ``slide_index + 1``).
        """
        source_slide = prs.slides[slide_index]

        # 1. Add a new blank slide with the same layout – this is appended at the end.
        new_slide = prs.slides.add_slide(source_slide.slide_layout)

        # 2. Replace the new slide's shape tree with a deep copy of the source's.
        source_spTree = source_slide.shapes._spTree
        new_spTree = new_slide.shapes._spTree

        for child in list(new_spTree):
            new_spTree.remove(child)
        for child in source_spTree:
            new_spTree.append(copy.deepcopy(child))

        # 3. Move the new slide from the end to slide_index + 1 by reordering
        #    the <p:sldIdLst> in the presentation XML.
        pml_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        try:
            presentation_elem = prs.part._element
            sldIdLst = presentation_elem.find(f'{{{pml_ns}}}sldIdLst')
            if sldIdLst is not None:
                sld_id_elems = list(sldIdLst)
                if sld_id_elems:
                    new_sldId = sld_id_elems[-1]          # just appended → at end
                    source_sldId = sld_id_elems[slide_index]  # source position
                    sldIdLst.remove(new_sldId)
                    source_sldId.addnext(new_sldId)
                    print(
                        f'  - Duplicated slide at index {slide_index}, '
                        f'inserted at index {slide_index + 1}',
                    )
        except Exception as e:
            print(f'  - Error repositioning duplicated slide: {e}')

        return prs.slides[slide_index + 1]

    # ------------------------------------------------------------------

    def _fill_slide_content(self, slide, slide_content: SlideContent):
        """Fill content into an existing slide"""
        slide_num = slide_content.slide_number
        config = self.template_config.get_slide_config(slide_num) if self.template_config else None

        # Strategy 1: Shape targeting (if shape_contents exists)
        if self._has_shape_contents(slide_content) and config and config.shape_targets:
            print(f'  - Using shape targeting for slide {slide_num}')
            self._fill_with_shape_targets(slide, slide_content, config)
            return

        # Strategy 2: Default filling
        self._fill_slide_default(slide, slide_content)

    def _has_shape_contents(self, slide_content: SlideContent) -> bool:
        """Check if slide_content has shape_contents attribute"""
        return bool(
            hasattr(slide_content, '__dict__') and
            'shape_contents' in slide_content.__dict__ and
            slide_content.shape_contents,
        )

    def _fill_with_shape_targets(self, slide, slide_content, config):
        """Fill content using shape targets from config"""
        shape_contents = getattr(slide_content, 'shape_contents', {})

        if not shape_contents:
            print('  - No shape_contents found in slide_content')
            return

        print(f'  - Found {len(shape_contents)} shape contents to fill')

        # Update title if exists
        self._update_slide_title(slide, slide_content.title)

        # Fill each shape target
        for target in config.shape_targets:
            self._fill_single_shape_target(slide, target, shape_contents, config)

    def _update_slide_title(self, slide, title: Optional[str]):
        """Update slide title if exists"""
        if slide.shapes.title and title:
            slide.shapes.title.text = title
            print(f'  - Updated title: {title}')

    def _fill_single_shape_target(self, slide, target, shape_contents: Dict[str, str], config=None):
        """Fill a single shape target with content"""
        content_key = target.content_key

        # Check if content exists for this key
        if content_key not in shape_contents:
            print(f"  - Content key '{content_key}' not found in markdown")
            return

        content = shape_contents[content_key]

        # Find shape by index or name
        shape = self._find_shape(slide, target)

        if not shape:
            print(f'  - Shape not found for target: {content_key}')
            return

        # Detect content type and fill accordingly
        if self._is_image_path(content):
            # Content is an image path
            self._replace_shape_with_image(slide, shape, content, target)
        elif target.fill_cols or self._is_table_data(content):
            # fill_cols always implies a table target (even single-row);
            # otherwise fall back to markdown table detection.
            self._fill_table(shape, content, target)
        elif shape.has_text_frame:
            # Content is text
            shape.text_frame.text = content
            self._apply_text_font(shape.text_frame)
            
            # Apply 10pt line spacing for slide 4
            if config and config.slide_number == 4:
                self._apply_line_spacing(shape.text_frame, Pt(10))
                print(f"  - Applied 10pt line spacing to slide 4 shape '{content_key}'")
            
            print(f"  Filled '{content_key}' into {shape.name} ({len(content)} chars)")
        else:
            print(f'  - Shape {shape.name} has no text frame and content is not an image or table')

    def _is_image_path(self, content: str) -> bool:
        """Check if content is an image file path or base64 data URI"""
        content = content.strip()
        # Check base64 data URI: data:image/...;base64,...
        if content.startswith('data:image/'):
            return True
        # Check if it's a file path with image extension
        image_extensions = ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.svg']
        return any(content.lower().endswith(ext) for ext in image_extensions)

    def _decode_base64_image(self, data_uri: str) -> Optional[str]:
        """Decode a base64 data URI to a temporary file and return the file path.

        Supports format: data:image/<ext>;base64,<data>

        Returns:
            Absolute path to the temporary file, or None on error.
        """
        try:
            # Parse header: data:image/png;base64,<data>
            header, encoded = data_uri.split(',', 1)
            # Extract mime type, e.g. image/png
            mime_part = header.split(':')[1].split(';')[0]  # e.g. "image/png"
            ext = '.' + mime_part.split('/')[1]  # e.g. ".png"
            # Handle special cases
            if ext == '.jpeg':
                ext = '.jpg'
            elif ext == '.svg+xml':
                ext = '.svg'

            image_bytes = base64.b64decode(encoded)
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
            tmp.write(image_bytes)
            tmp.flush()
            tmp.close()
            print(f'  - Decoded base64 image to temp file: {tmp.name} ({len(image_bytes)} bytes)')
            return tmp.name
        except Exception as e:
            print(f'  - Failed to decode base64 image: {e}')
            return None

    def _replace_shape_with_image(self, slide, old_shape, image_path: str, target):
        """Replace a shape with an image (supports file path or base64 data URI)"""
        tmp_file: Optional[str] = None
        try:
            # Get shape position and size
            left = old_shape.left
            top = old_shape.top
            width = old_shape.width
            height = old_shape.height

            image_path = image_path.strip()

            # Handle base64 data URI
            if image_path.startswith('data:image/'):
                tmp_file = self._decode_base64_image(image_path)
                if not tmp_file:
                    print(f"  - Cannot decode base64 image for shape '{target.content_key}'")
                    return
                image_path = tmp_file
            else:
                # Resolve file path
                if not Path(image_path).is_absolute():
                    rel_path = Path(image_path)
                    if rel_path.exists():
                        image_path = str(rel_path)
                    elif self._output_dir:
                        resolved = self._output_dir / image_path
                        if resolved.exists():
                            image_path = str(resolved)
                    else:
                        outputs_path = Path('outputs') / image_path
                        if outputs_path.exists():
                            image_path = str(outputs_path)

                # Check if image exists
                if not Path(image_path).exists():
                    print(f'  - Image not found: {image_path}')
                    return

            # Add new image at same position (overlays the placeholder shape)
            try:
                slide.shapes.add_picture(
                    str(image_path),
                    left, top,
                    width=width,
                    height=height,
                )
                print(f"  Replaced shape '{target.content_key}' with image: {Path(image_path).name}")
            except Exception as e:
                print(f'  - Error adding image: {e}')

        except Exception as e:
            print(f'  - Error replacing shape with image: {e}')
        finally:
            # Clean up temp file if created from base64
            if tmp_file and Path(tmp_file).exists():
                try:
                    Path(tmp_file).unlink()
                    print(f'  - Cleaned up temp file: {tmp_file}')
                except Exception:
                    pass

    def _is_table_data(self, content: str) -> bool:
        """Check if content is table data (markdown table format).

        Accepts two formats:
        1. Standard markdown table with header + separator:
               | col1 | col2 |
               |------|------|
               | val1 | val2 |
        2. No-header row-only table (used for fill_cols targeting):
               | val1 |
               | val2 |
               | val3 |
        """
        content = content.strip()
        lines = [line.strip() for line in content.split('\n') if line.strip()]

        if len(lines) < 2:
            return False

        # All lines must contain at least one pipe
        if not all('|' in line for line in lines):
            return False

        # Format 1: standard table with separator line
        for line in lines[1:3]:
            if all(c in '|-: \t' for c in line):
                return True

        # Format 2: no-separator — every line starts and ends with | (row-only table)
        if all(line.startswith('|') and line.endswith('|') for line in lines):
            return True

        return False

    def _fill_table(self, shape, table_data: str, target):
        """Fill an existing table shape with data from markdown table format"""
        try:
            # Check if shape has table
            if not shape.has_table:
                print(f"  - Shape '{target.content_key}' is not a table")
                return

            table = shape.table

            # Parse markdown table
            lines = [line.strip() for line in table_data.strip().split('\n') if line.strip()]

            # Remove separator line (e.g., |---|---|)
            table_lines = []
            for line in lines:
                # Skip separator lines
                if all(c in '|-: \t' for c in line):
                    continue
                table_lines.append(line)

            if not table_lines:
                print('  - No valid table data found')
                return

            # Parse cells from each line
            rows_data = []
            for line in table_lines:
                # Split by | and clean up
                cells = [cell.strip() for cell in line.split('|')]
                # Remove empty first/last cells (from leading/trailing |)
                cells = [cell for cell in cells if cell]
                if cells:
                    rows_data.append(cells)

            if not rows_data:
                print('  - No valid rows found in table data')
                return

            # Get current table dimensions
            current_rows = len(table.rows)
            current_cols = len(table.columns)
            needed_rows = len(rows_data)
            needed_cols = len(rows_data[0]) if rows_data else 0

            # Adjust table size if needed
            if needed_rows > current_rows:
                # Add rows if we need more
                rows_to_add = needed_rows - current_rows
                print(f'  - Adding {rows_to_add} rows to table (current: {current_rows}, needed: {needed_rows})')
                self._add_table_rows(table, rows_to_add)
                # After adding rows, update current_rows to reflect the new count.
                # Note: python-pptx's table.rows collection is live, so this should
                # reflect the new size immediately, but we recalculate for safety.
                current_rows = len(table.rows)
            elif needed_rows < current_rows:
                # Fewer rows needed - remove the extra rows from the XML so they
                # don't appear as empty rows in the rendered slide.
                extra_rows = current_rows - needed_rows
                print(
                    '  - Table has %s rows but only %s rows of data '
                    '(removing %s unused rows)',
                    current_rows,
                    needed_rows,
                    extra_rows,
                )
                self._remove_table_rows(table, needed_rows)
                current_rows = len(table.rows)

            if needed_cols > current_cols:
                print(f'  - Data has {needed_cols} columns but table only has {current_cols} columns')
                needed_cols = current_cols  # Limit to available columns

            # Fill table cells
            filled_count = 0
            for row_idx, row_data in enumerate(rows_data):
                # After adding rows via XML, we can fill up to needed_rows
                # Use try-except to handle any edge cases where table structure might not match
                try:
                    # If fill_cols is specified, map row_data positionally onto those columns;
                    # columns not listed are left untouched (preserving template content).
                    if target.fill_cols:
                        col_pairs = list(zip(target.fill_cols, row_data))
                    else:
                        col_pairs = list(enumerate(row_data))

                    for col_idx, cell_value in col_pairs:
                        if col_idx >= current_cols:
                            break  # Skip extra columns in data

                        cell = table.cell(row_idx, col_idx)
                        cell.text = cell_value
                        self._apply_table_cell_font(cell)
                        filled_count += 1
                except (IndexError, AttributeError) as e:
                    print(f'  - Cannot fill row {row_idx}: {e}')
                    break  # Stop filling if we hit an error

            print(f"  Filled '{target.content_key}' table: {len(rows_data)} rows x {needed_cols} cols ({filled_count} cells)")

        except Exception as e:
            print(f'  - Error filling table: {e}')

    def _add_table_rows(self, table, num_rows: int):
        """Add rows to an existing table by manipulating XML

        python-pptx uses lxml internally, so we can manipulate the XML directly
        to add new rows by cloning the structure of the last row.
        """
        try:
            from lxml import etree

            # Get the table XML element (this is an lxml element)
            tbl = table._tbl

            # Get the number of existing rows
            current_row_count = len(table.rows)

            # Get the last row to use as template for new rows
            if current_row_count == 0:
                print('  - Cannot add rows: table has no existing rows to use as template')
                return

            # Use positive index instead of negative index (table.rows doesn't support negative indexing)
            last_row = table.rows[current_row_count - 1]
            last_row_xml = last_row._tr

            # Define namespace for XPath queries
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

            # Add new rows by cloning the last row
            for i in range(num_rows):
                # Create a deep copy of the last row by serializing and parsing
                new_row_xml = etree.fromstring(etree.tostring(last_row_xml))

                # Clear text content in all cells of the new row
                for tc in new_row_xml.xpath('.//a:tc', namespaces=ns):
                    txBody = tc.find('.//a:txBody', namespaces=ns)
                    if txBody is not None:
                        # Get all paragraphs
                        paragraphs = txBody.findall('.//a:p', namespaces=ns)
                        if paragraphs:
                            # Clear text in first paragraph
                            first_p = paragraphs[0]
                            for t_elem in first_p.xpath('.//a:t', namespaces=ns):
                                t_elem.text = ''
                            # Remove other paragraphs (keep structure but remove extra content)
                            for p in paragraphs[1:]:
                                txBody.remove(p)

                # Append the new row to the table
                tbl.append(new_row_xml)

            print(f'  - Successfully added {num_rows} rows to table')

        except ImportError:
            print('  - lxml not available. python-pptx requires lxml, please install it: pip install lxml')
        except Exception as e:
            print(f'  - Error adding rows to table: {e}')
            import traceback
            print(f'  - Traceback: {traceback.format_exc()}')

    def _remove_table_rows(self, table, keep_rows: int):
        """Remove all rows after index *keep_rows* from the table by manipulating XML.

        Args:
            table: python-pptx Table object
            keep_rows: Number of rows to keep (rows at index 0 … keep_rows-1 are retained).
        """
        try:
            from lxml import etree  # noqa: F401 (just to confirm availability)

            tbl = table._tbl
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            tr_elements = tbl.findall('a:tr', namespaces=ns)

            rows_to_remove = tr_elements[keep_rows:]
            for tr in rows_to_remove:
                tbl.remove(tr)

            print(
                f'  - Removed {len(rows_to_remove)} extra row(s); '
                f'table now has {keep_rows} row(s)',
            )
        except Exception as e:
            print(f'  - Error removing table rows: {e}')

    def _find_shape(self, slide, target):
        """Find shape by index or name.

        When only ``shape_name`` is provided (no ``shape_index``), the search
        is performed recursively through GROUP shapes so that shapes nested
        inside groups can be targeted by name.
        """
        try:
            # Try by index first (faster, top-level only)
            if target.shape_index is not None and target.shape_index < len(slide.shapes):
                shape = slide.shapes[target.shape_index]
                print(f'  - Found shape by index {target.shape_index}: {shape.name}')
                return shape

            # Fallback to recursive name search (traverses group children)
            if target.shape_name:
                found = self._find_shape_recursive(slide.shapes, target.shape_name)
                if found:
                    print(f'  - Found shape by name (recursive): {target.shape_name}')
                    return found
        except Exception as e:
            print(f'  - Error finding shape: {e}')

        return None

    def _find_shape_recursive(self, shapes, name: str):
        """Recursively search for a shape by name, traversing GROUP children.

        Args:
            shapes: A shape collection (``slide.shapes`` or ``group_shape.shapes``).
            name: The exact shape name to search for.

        Returns:
            The matching shape, or ``None`` if not found.
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            for shape in shapes:
                if shape.name == name:
                    return shape
                # Recurse into group shapes
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        found = self._find_shape_recursive(shape.shapes, name)
                        if found:
                            return found
                except AttributeError:
                    pass
        except Exception as e:
            print(f'  - _find_shape_recursive error: {e}')
        return None

    def _fill_slide_default(self, slide, slide_content: SlideContent):
        """Default filling strategy (original logic)"""

        # Update title
        self._update_slide_title(slide, slide_content.title)

        # Get available text containers
        placeholders = self._get_text_placeholders(slide)
        text_shapes = self._get_text_shapes(slide)

        # Fill content items
        placeholder_idx = 0
        for content_item in slide_content.content:
            if content_item.type == ContentType.TEXT:
                placeholder_idx = self._fill_text_content(
                    slide, content_item, placeholders, text_shapes, placeholder_idx,
                )

            elif content_item.type == ContentType.LIST:
                placeholder_idx = self._fill_list_content(
                    slide, content_item, placeholders, placeholder_idx,
                )

            elif content_item.type == ContentType.TABLE:
                self._add_table(slide, content_item, Inches(1.0), Inches(2.5), Inches(8.0))

            elif content_item.type == ContentType.IMAGE:
                self._add_image(slide, content_item, Inches(1.0), Inches(2.5), Inches(8.0))

    def _get_text_placeholders(self, slide):
        """Get all text placeholders in slide (excluding title)"""
        return [
            s for s in slide.shapes
            if s.is_placeholder and s.has_text_frame and s != slide.shapes.title
        ]

    def _get_text_shapes(self, slide):
        """Get all text shapes in slide (excluding placeholders and title)"""
        return [
            s for s in slide.shapes
            if s.has_text_frame and not s.is_placeholder and s != slide.shapes.title
        ]

    def _fill_text_content(self, slide, content_item, placeholders, text_shapes, placeholder_idx):
        """Fill text content into placeholders or text shapes"""
        if placeholder_idx < len(placeholders):
            placeholders[placeholder_idx].text_frame.text = content_item.data
            self._apply_text_font(placeholders[placeholder_idx].text_frame)
            print(f'  - Filled placeholder {placeholder_idx}: {content_item.data[:50]}...')
            return placeholder_idx + 1

        elif text_shapes:
            text_shapes[0].text_frame.text = content_item.data
            self._apply_text_font(text_shapes[0].text_frame)
            text_shapes.pop(0)
            return placeholder_idx

        else:
            # Add new textbox if no placeholders left
            self._add_content_at_position(slide, content_item, Inches(1.0), Inches(2.0), Inches(8.0))
            return placeholder_idx

    def _fill_list_content(self, slide, content_item, placeholders, placeholder_idx):
        """Fill list content into placeholder"""
        if placeholder_idx < len(placeholders):
            text_frame = placeholders[placeholder_idx].text_frame
            text_frame.clear()
            for i, item in enumerate(content_item.data):
                if i > 0:
                    text_frame.add_paragraph()
                p = text_frame.paragraphs[i]
                p.text = f'• {item}'
                p.level = 0
            self._apply_text_font(text_frame)
            print(f'  - Filled list with {len(content_item.data)} items')
            return placeholder_idx + 1
        else:
            # Add new list
            top = Inches(2.5 + (placeholder_idx * 0.5))
            self._add_content_at_position(slide, content_item, Inches(1.0), top, Inches(8.0))
            return placeholder_idx

    def _add_content_at_position(self, slide, content_item: ContentItem, left, top, width):
        """Add content item at specific position"""
        if content_item.type == ContentType.TEXT:
            height = Inches(0.5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            textbox.text_frame.text = content_item.data
            self._apply_text_font(textbox.text_frame)

        elif content_item.type == ContentType.LIST:
            height = Inches(0.3) * len(content_item.data)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame

            for i, item in enumerate(content_item.data):
                if i > 0:
                    text_frame.add_paragraph()
                p = text_frame.paragraphs[i]
                p.text = f'• {item}'
            self._apply_text_font(text_frame)

    def _add_image(self, slide, content_item: ContentItem, left, top, width):
        """Add image to slide"""
        try:
            image_data = content_item.data
            image_path = self._resolve_image_path(image_data)

            if Path(image_path).exists():
                slide.shapes.add_picture(str(image_path), left, top, width=width)
                self._add_image_caption(slide, image_data, left, top, width)
            else:
                print(f'Image not found: {image_path}')
                self._add_placeholder_text(slide, f'[Image: {image_path}]', left, top, width)

        except Exception as e:
            print(f'Error adding image: {e}')

    def _resolve_image_path(self, image_data) -> str:
        """Resolve image path from data"""
        image_path = image_data.get('path') if isinstance(image_data, dict) else image_data

        if image_path and not Path(str(image_path)).is_absolute():
            # Try relative to CWD first
            rel_path = Path(str(image_path))
            if not rel_path.exists():
                # Try outputs directory
                image_path = str(Path('outputs') / str(image_path))

        return str(image_path) if image_path else ''

    def _add_image_caption(self, slide, image_data, left, top, width):
        """Add caption below image if exists"""
        if isinstance(image_data, dict) and image_data.get('caption'):
            caption_top = top + Inches(3.0)
            caption_box = slide.shapes.add_textbox(left, caption_top, width, Inches(0.3))
            caption_box.text_frame.text = image_data['caption']
            for paragraph in caption_box.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'Noto Sans JP'
                    run.font.size = Pt(10)
                    run.font.italic = True

    def _apply_text_font(self, text_frame, font_name: str = 'Noto Sans JP', font_size: int = 10):
        """Apply font and size to all text in a text frame"""
        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font_name
                    run.font.size = Pt(font_size)
        except Exception as e:
            print(f'  - Error applying font to text: {e}')

    def _apply_line_spacing(self, text_frame, spacing_pt):
        """Apply line spacing (space after) to all paragraphs in a text frame
        
        Args:
            text_frame: The text frame to apply spacing to
            spacing_pt: Spacing in points (e.g., Pt(10) for 10pt)
        """
        try:
            for paragraph in text_frame.paragraphs:
                if paragraph.text.strip():  # skip empty paragraphs (from blank lines)
                    paragraph.space_after = spacing_pt
        except Exception as e:
            print(f'  - Error applying line spacing: {e}')

    def _apply_table_cell_font(self, cell, font_name: str = 'Noto Sans JP', font_size: int = 8):
        """Apply font and size to all text in a table cell"""
        try:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font_name
                    run.font.size = Pt(font_size)
        except Exception as e:
            print(f'  - Error applying font to table cell: {e}')

    def _add_placeholder_text(self, slide, text: str, left, top, width):
        """Add placeholder text box"""
        textbox = slide.shapes.add_textbox(left, top, width, Inches(1))
        textbox.text_frame.text = text
        self._apply_text_font(textbox.text_frame)

    def _add_table(self, slide, content_item: ContentItem, left, top, width):
        """Add table to slide"""
        try:
            table_data = content_item.data
            headers = table_data.get('headers', [])
            rows = table_data.get('rows', [])

            if not headers or not rows:
                print('Table has no headers or rows')
                return

            # Create table
            row_count = len(rows) + 1  # +1 for header
            col_count = len(headers)
            height = Inches(0.3) * row_count

            table_shape = slide.shapes.add_table(row_count, col_count, left, top, width, height)
            table = table_shape.table

            # Fill headers
            self._fill_table_headers(table, headers)

            # Fill data rows
            self._fill_table_rows(table, rows)

        except Exception as e:
            print(f'Error adding table: {e}')

    def _fill_table_headers(self, table, headers: List[str]):
        """Fill table header row"""
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            # Apply font and make header bold
            self._apply_table_cell_font(cell)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True

    def _fill_table_rows(self, table, rows: List[List[str]]):
        """Fill table data rows"""
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)  # +1 to skip header
                cell.text = str(cell_data)
                self._apply_table_cell_font(cell)
